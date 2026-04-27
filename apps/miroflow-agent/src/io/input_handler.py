# Copyright (c) 2025 MiroMind
# This source code is licensed under the Apache 2.0 License.

"""
Input handler module for processing various file types.

This module provides functions for:
- Processing task inputs with associated files
- Converting documents (PDF, DOCX, PPTX, XLSX) to markdown
- Generating captions for images, audio, and video files
- Extracting task-relevant information from media files

Supported file formats:
- Documents: PDF, DOCX, DOC, PPTX, PPT, XLSX, XLS, HTML
- Images: JPG, JPEG, PNG, GIF, WEBP
- Audio: WAV, MP3, M4A
- Video: MP4, MOV, AVI, MKV, WEBM
- Data: JSON, JSONLD, CSV, YAML, TOML
- Code: PY, SH, MD, TXT
- Archives: ZIP
"""

import base64
import html
import json
import os
import re
import shutil
import tempfile
import traceback
from typing import Any, Tuple, Union
from urllib.parse import quote, unquote, urlparse, urlunparse

import mammoth
import requests
import markdownify
import openpyxl
import pdfminer
import pdfminer.high_level
import pptx
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from markitdown import MarkItDown
from openai import OpenAI
from openpyxl.utils import get_column_letter

# Ensure .env file is loaded
load_dotenv()

# Local vision endpoint (same as tool-vqa-os)
VISION_API_KEY = os.environ.get("VISION_API_KEY", "not-needed")
VISION_BASE_URL = os.environ.get("VISION_BASE_URL", "http://localhost:8001/v1/chat/completions")
VISION_MODEL_NAME = os.environ.get("VISION_MODEL_NAME", "qwen3.5")

# File extension constants for different media types
IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "gif", "webp"}
AUDIO_EXTENSIONS = {"wav", "mp3", "m4a"}
VIDEO_EXTENSIONS = {"mp4", "mov", "avi", "mkv", "webm"}
MEDIA_EXTENSIONS = IMAGE_EXTENSIONS | AUDIO_EXTENSIONS | VIDEO_EXTENSIONS
# Extensions that should skip MarkItDown fallback processing
SKIP_MARKITDOWN_EXTENSIONS = MEDIA_EXTENSIONS | {"pdb"}


def _generate_image_caption(image_path: str) -> str:
    """
    Generate a caption for an image using the local vision LLM (qwen3.5).

    Args:
        image_path: Path to the image file

    Returns:
        Caption string, or error message if failed
    """
    try:
        # Read and encode image
        with open(image_path, "rb") as image_file:
            image_data = base64.b64encode(image_file.read()).decode("utf-8")

        # Guess MIME type
        _, ext = os.path.splitext(image_path)
        ext = ext.lower()
        mime_type = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".gif": "image/gif",
            ".webp": "image/webp",
        }.get(ext, "image/jpeg")

        payload = {
            "model": VISION_MODEL_NAME,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Please provide a detailed description of this image. Include key objects, people, text, colors, and any other relevant details.",
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{image_data}"
                            },
                        },
                    ],
                }
            ],
            "max_tokens": 2048,
            "temperature": 0,
        }

        headers = {
            "Authorization": f"Bearer {VISION_API_KEY}",
            "Content-Type": "application/json",
        }

        response = requests.post(VISION_BASE_URL, json=payload, headers=headers, timeout=120)
        response.raise_for_status()
        content = response.json()["choices"][0]["message"]["content"]
        return content if content else "[Caption unavailable: Empty response]"

    except Exception as e:
        return f"[Caption generation failed: {str(e)}]"


def _generate_audio_caption(audio_path: str) -> str:
    """
    Generate a caption for an audio file using OpenAI's audio transcription.

    Args:
        audio_path: Path to the audio file

    Returns:
        Caption string (transcription), or error message if failed
    """
    try:
        OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
        OPENAI_BASE_URL = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")

        if not OPENAI_API_KEY:
            return "[Caption unavailable: OPENAI_API_KEY not set]"

        client = OpenAI(api_key=OPENAI_API_KEY, base_url=OPENAI_BASE_URL)

        # Transcribe audio
        with open(audio_path, "rb") as audio_file:
            transcription = client.audio.transcriptions.create(
                model="gpt-4o-transcribe", file=audio_file
            )

        text = transcription.text
        return text if text else "[Transcription unavailable: Empty response]"

    except Exception as e:
        return f"[Caption generation failed: {str(e)}]"


def _generate_video_caption(video_path: str) -> str:
    """
    Generate a caption for a video using the local vision LLM (qwen3.5).

    Args:
        video_path: Path to the video file

    Returns:
        Caption string, or error message if failed
    """
    try:
        # Read and encode video
        with open(video_path, "rb") as video_file:
            video_data = base64.b64encode(video_file.read()).decode("utf-8")

        # Guess MIME type
        _, ext = os.path.splitext(video_path)
        ext = ext.lower()
        mime_type = {
            ".mp4": "video/mp4",
            ".mov": "video/quicktime",
            ".avi": "video/x-msvideo",
            ".mkv": "video/x-matroska",
            ".webm": "video/webm",
        }.get(ext, "video/mp4")

        payload = {
            "model": VISION_MODEL_NAME,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Please provide a detailed description of this video. Include key events, people, objects, actions, audio information, and any text visible in the video.",
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{video_data}"
                            },
                        },
                    ],
                }
            ],
            "max_tokens": 2048,
            "temperature": 0,
        }

        headers = {
            "Authorization": f"Bearer {VISION_API_KEY}",
            "Content-Type": "application/json",
        }

        response = requests.post(VISION_BASE_URL, json=payload, headers=headers, timeout=120)
        response.raise_for_status()
        content = response.json()["choices"][0]["message"]["content"]
        return content if content else "[Caption unavailable: Empty response]"

    except Exception as e:
        return f"[Caption generation failed: {str(e)}]"


def _extract_task_relevant_info_from_image(
    image_path: str, task_description: str
) -> str:
    """
    Extract task-relevant information directly from an image based on the task description.

    Args:
        image_path: Path to the image file
        task_description: The user's task description

    Returns:
        Extracted relevant information, or empty string if extraction fails
    """
    try:
        # Read and encode image
        with open(image_path, "rb") as image_file:
            image_data = base64.b64encode(image_file.read()).decode("utf-8")

        # Guess MIME type
        _, ext = os.path.splitext(image_path)
        ext = ext.lower()
        mime_type = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".gif": "image/gif",
            ".webp": "image/webp",
        }.get(ext, "image/jpeg")

        payload = {
            "model": VISION_MODEL_NAME,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"""Based on the following task, analyze this image and extract only the information that is directly relevant to completing the task.

Task: {task_description}

Please provide a concise summary of the relevant information from the image that would help in completing this task. Focus only on what's pertinent to the task. If nothing is particularly relevant, state "No specific task-relevant details identified in the image." Keep the response brief and focused.""",
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{image_data}"
                            },
                        },
                    ],
                }
            ],
            "max_tokens": 1024,
            "temperature": 0,
        }

        headers = {
            "Authorization": f"Bearer {VISION_API_KEY}",
            "Content-Type": "application/json",
        }

        response = requests.post(VISION_BASE_URL, json=payload, headers=headers, timeout=120)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"].strip()

    except Exception as e:
        print(f"Warning: Failed to extract task-relevant info from image: {str(e)}")
        return ""


def _extract_task_relevant_info_from_audio(
    audio_path: str, task_description: str
) -> str:
    """
    Extract task-relevant information directly from an audio file based on the task description.

    Args:
        audio_path: Path to the audio file
        task_description: The user's task description

    Returns:
        Extracted relevant information, or empty string if extraction fails
    """
    try:
        OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
        OPENAI_BASE_URL = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")

        if not OPENAI_API_KEY:
            return ""

        client = OpenAI(api_key=OPENAI_API_KEY, base_url=OPENAI_BASE_URL)

        # Read and encode audio file
        with open(audio_path, "rb") as audio_file:
            audio_data = base64.b64encode(audio_file.read()).decode("utf-8")

        # Detect audio format
        _, ext = os.path.splitext(audio_path)
        ext = ext.lower()
        audio_format = {
            ".mp3": "mp3",
            ".wav": "wav",
            ".m4a": "m4a",
        }.get(ext, "mp3")

        # Use gpt-4o-audio-preview for direct audio question answering
        text_prompt = f"""Based on the following task, analyze this audio and extract only the information that is directly relevant to completing the task.

Task: {task_description}

Please provide a concise summary of the relevant information from the audio that would help in completing this task. Focus only on what's pertinent to the task. If nothing is particularly relevant, state "No specific task-relevant details identified in the audio." Keep the response brief and focused."""

        response = client.chat.completions.create(
            model="gpt-4o-audio-preview",
            messages=[
                {
                    "role": "system",
                    "content": "You are a helpful assistant specializing in audio analysis.",
                },
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": text_prompt},
                        {
                            "type": "input_audio",
                            "input_audio": {
                                "data": audio_data,
                                "format": audio_format,
                            },
                        },
                    ],
                },
            ],
            max_tokens=1024,
            temperature=0,
        )

        return response.choices[0].message.content.strip()

    except Exception as e:
        print(f"Warning: Failed to extract task-relevant info from audio: {str(e)}")
        return ""


def _extract_task_relevant_info_from_video(
    video_path: str, task_description: str
) -> str:
    """
    Extract task-relevant information directly from a video based on the task description.

    Args:
        video_path: Path to the video file
        task_description: The user's task description

    Returns:
        Extracted relevant information, or empty string if extraction fails
    """
    try:
        # Read and encode video
        with open(video_path, "rb") as video_file:
            video_data = base64.b64encode(video_file.read()).decode("utf-8")

        # Guess MIME type
        _, ext = os.path.splitext(video_path)
        ext = ext.lower()
        mime_type = {
            ".mp4": "video/mp4",
            ".mov": "video/quicktime",
            ".avi": "video/x-msvideo",
            ".mkv": "video/x-matroska",
            ".webm": "video/webm",
        }.get(ext, "video/mp4")

        payload = {
            "model": VISION_MODEL_NAME,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"""Based on the following task, analyze this video and extract only the information that is directly relevant to completing the task.

Task: {task_description}

Please provide a concise summary of the relevant information from the video that would help in completing this task. Focus only on what's pertinent to the task. If nothing is particularly relevant, state "No specific task-relevant details identified in the video." Keep the response brief and focused.""",
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{video_data}"
                            },
                        },
                    ],
                }
            ],
            "max_tokens": 1024,
            "temperature": 0,
        }

        headers = {
            "Authorization": f"Bearer {VISION_API_KEY}",
            "Content-Type": "application/json",
        }

        response = requests.post(VISION_BASE_URL, json=payload, headers=headers, timeout=120)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"].strip()

    except Exception as e:
        print(f"Warning: Failed to extract task-relevant info from video: {str(e)}")
        return ""


def process_input(task_description: str, task_file_name: str) -> Tuple[str, str]:
    """
    Process user input and associated files.

    Extracts content from the task file (if provided) and appends it to the
    task description in a format suitable for the LLM.

    Args:
        task_description: The original task description
        task_file_name: Path to an associated file, or empty string if none

    Returns:
        Tuple of (updated_task_description, updated_task_description)
        Both values are the same - the task description with file content appended
    """
    updated_task_description = task_description
    file_content_section = ""  # Collect file content to append at the end

    if task_file_name:
        try:
            file_extension = task_file_name.rsplit(".", maxsplit=1)[-1].lower()
            parsing_result = None

            if file_extension in IMAGE_EXTENSIONS:
                # Generate unconditional image caption
                caption = _generate_image_caption(task_file_name)

                # Extract task-relevant information directly from the image
                relevant_info = _extract_task_relevant_info_from_image(
                    task_file_name, task_description
                )

                # Format as Markdown
                file_content_section += f"\n\nNote: An image file '{task_file_name}' is associated with this task. The content has been extracted as a detailed caption below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## Image Content\nFile: {task_file_name}\n\n"
                file_content_section += f"> {caption}\n\n"

                if relevant_info:
                    file_content_section += "Task-Relevant Information:\n\n"
                    file_content_section += f"{relevant_info}\n\n"

            elif file_extension == "py":
                # Python files - read directly
                with open(task_file_name, "r", encoding="utf-8") as f:
                    parsing_result = DocumentConverterResult(
                        title=None, text_content=f.read()
                    )
                file_content_section += f"\n\nNote: A Python file '{task_file_name}' is associated with this task. The content has been extracted as text below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## Python File\nFile: {task_file_name}\n\n"

            elif file_extension in ["txt", "md", "sh", "yaml", "yml", "toml", "csv"]:
                # Text-based files - read directly
                with open(task_file_name, "r", encoding="utf-8") as f:
                    parsing_result = DocumentConverterResult(
                        title=None, text_content=f.read()
                    )
                file_type_name = {
                    "txt": "Text",
                    "md": "Markdown",
                    "sh": "Shell Script",
                    "yaml": "YAML",
                    "yml": "YAML",
                    "toml": "TOML",
                    "csv": "CSV",
                }.get(file_extension, "Text")
                file_content_section += f"\n\nNote: A {file_type_name.lower()} file '{task_file_name}' is associated with this task. The content has been extracted as text below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += (
                    f"## {file_type_name} File\nFile: {task_file_name}\n\n"
                )

            elif file_extension in ["jsonld", "json"]:
                with open(task_file_name, "r", encoding="utf-8") as f:
                    parsing_result = DocumentConverterResult(
                        title=None,
                        text_content=json.dumps(
                            json.load(f), ensure_ascii=False, indent=2
                        ),
                    )
                file_content_section += f"\n\nNote: A JSON file '{task_file_name}' is associated with this task. The content has been extracted as JSON format below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## JSON File\nFile: {task_file_name}\n\n"

            elif file_extension in ["xlsx", "xls"]:
                parsing_result = XlsxConverter(local_path=task_file_name)
                file_content_section += f"\n\nNote: An Excel file '{task_file_name}' is associated with this task. The content has been extracted as a markdown table below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## Excel File\nFile: {task_file_name}\n\n"

            elif file_extension == "pdf":
                parsing_result = DocumentConverterResult(
                    title=None,
                    text_content=pdfminer.high_level.extract_text(task_file_name),
                )
                file_content_section += f"\n\nNote: A PDF file '{task_file_name}' is associated with this task. The content has been extracted as text below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## PDF File\nFile: {task_file_name}\n\n"

            elif file_extension in ["docx", "doc"]:
                parsing_result = DocxConverter(local_path=task_file_name)
                file_content_section += f"\n\nNote: A Word document '{task_file_name}' is associated with this task. The content has been extracted as markdown below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## Word Document\nFile: {task_file_name}\n\n"

            elif file_extension in ["html", "htm"]:
                parsing_result = HtmlConverter(local_path=task_file_name)
                file_content_section += f"\n\nNote: An HTML file '{task_file_name}' is associated with this task. The content has been extracted as markdown below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## HTML File\nFile: {task_file_name}\n\n"

            elif file_extension in ["pptx", "ppt"]:
                parsing_result = PptxConverter(local_path=task_file_name)
                file_content_section += f"\n\nNote: A PowerPoint presentation '{task_file_name}' is associated with this task. The content has been extracted as markdown below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += (
                    f"## PowerPoint Presentation\nFile: {task_file_name}\n\n"
                )

            elif file_extension in AUDIO_EXTENSIONS:
                # Generate unconditional audio transcription
                caption = _generate_audio_caption(task_file_name)

                # Extract task-relevant information directly from the audio
                relevant_info = _extract_task_relevant_info_from_audio(
                    task_file_name, task_description
                )

                # Format as Markdown
                file_content_section += f"\n\nNote: An audio file '{task_file_name}' is associated with this task. The content has been extracted as a transcription below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## Audio Content\nFile: {task_file_name}\n\n"
                file_content_section += f"> {caption}\n\n"

                if relevant_info:
                    file_content_section += "Task-Relevant Information:\n\n"
                    file_content_section += f"{relevant_info}\n\n"

            elif file_extension in VIDEO_EXTENSIONS:
                # Generate unconditional video caption
                caption = _generate_video_caption(task_file_name)

                # Extract task-relevant information directly from the video
                relevant_info = _extract_task_relevant_info_from_video(
                    task_file_name, task_description
                )

                # Format as Markdown
                file_content_section += f"\n\nNote: A video file '{task_file_name}' is associated with this task. The content has been extracted as a detailed caption below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## Video Content\nFile: {task_file_name}\n\n"
                file_content_section += f"> {caption}\n\n"

                if relevant_info:
                    file_content_section += "Task-Relevant Information:\n\n"
                    file_content_section += f"{relevant_info}\n\n"

            elif file_extension in ["zip"]:
                parsing_result = ZipConverter(local_path=task_file_name)
                file_content_section += f"\n\nNote: A ZIP archive '{task_file_name}' is associated with this task. The content has been extracted as file list and contents below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                file_content_section += f"## ZIP Archive\nFile: {task_file_name}\n\n"

            elif file_extension == "pdb":
                # PDB files (protein database) - only add note
                file_content_section += f"\n\nNote: A PDB file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"

            else:
                # For other file types, let MarkItDown try to handle it
                pass  # MarkItDown will be tried below

            #### markitdown process - ONLY if no specialized converter handled it ####
            if parsing_result is None:
                try:
                    if file_extension not in SKIP_MARKITDOWN_EXTENSIONS:
                        md = MarkItDown(enable_plugins=True)
                        parsing_result = md.convert(task_file_name)
                        print(
                            f"Info: Used MarkItDown as fallback to process file {task_file_name}"
                        )
                        # Add prompt for files processed by MarkItDown
                        file_content_section += f"\n\nNote: A file '{task_file_name}' is associated with this task. The content has been extracted as markdown below. You may use available tools to process its content if necessary. If you need to further process this file in the sandbox, please upload it to the sandbox first.\n\n"
                        file_content_section += (
                            f"## File Content\nFile: {task_file_name}\n\n"
                        )
                except Exception as e:
                    print(
                        f"Warning: MarkItDown failed to process {task_file_name}: {e}"
                    )
                    pass
            ############################

            # Collect the content and title (if has) to append later
            if getattr(parsing_result, "title", None):
                file_content_section += "Title:\n\n{}\n\n".format(parsing_result.title)
                file_content_section += "Content:\n\n```\n{}\n```\n".format(
                    parsing_result.text_content
                )
            elif getattr(parsing_result, "text_content", None):
                content = parsing_result.text_content
                max_len = 200_000  # Limit the length of results returned to LLM
                if len(content) > max_len:
                    content = content[:max_len] + "\n... [File truncated]"
                file_content_section += "```\n{}\n```\n".format(content)
            else:
                pass  # for image, audio, video files that already have their content formatted

        except FileNotFoundError:
            print(f"Error: File not found {task_file_name}")
            file_content_section += (
                f"\nWarning: The specified file '{task_file_name}' was not found."
            )
        except Exception as e:
            print(f"Error: Error processing file {task_file_name}: {e}")
            traceback.print_exc()
            file_content_section += f"\nWarning: There was an error processing the file '{task_file_name}': {str(e)}"

    # output format requirement
    updated_task_description += "\nYou should follow the format instruction in the request strictly and wrap the final answer in \\boxed{}."

    # Append file content at the end
    updated_task_description += file_content_section
    updated_task_description = updated_task_description.strip()

    return updated_task_description, updated_task_description


class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in [
                    "http",
                    "https",
                    "file",
                ]:  # type: ignore
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(
                    parsed_url._replace(path=quote(unquote(parsed_url.path)))
                )  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return "%s%s%s" % (prefix, text, suffix)

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return (
            "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix)
            if href
            else text
        )

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if (
            convert_as_inline
            and el.parent.name not in self.options["keep_inline_images_in"]
        ):
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)  # type: ignore


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


def convert_html_to_md(html_content):
    """
    Placeholder for HTML to Markdown conversion function
    In the original class, this would call self._convert()
    """
    soup = BeautifulSoup(html_content, "html.parser")
    for script in soup(["script", "style"]):
        script.extract()

    # Print only the main content
    body_elm = soup.find("body")
    webpage_text = ""
    if body_elm:
        webpage_text = _CustomMarkdownify().convert_soup(body_elm)
    else:
        webpage_text = _CustomMarkdownify().convert_soup(soup)

    assert isinstance(webpage_text, str)

    return DocumentConverterResult(
        title=None if soup.title is None else soup.title.string,
        text_content=webpage_text,
    )


def HtmlConverter(local_path: str):
    """
    Convert an HTML file to Markdown format.

    Args:
        local_path: Path to the HTML file to convert.

    Returns:
        DocumentConverterResult containing the converted Markdown text.
    """
    with open(local_path, "rt", encoding="utf-8") as fh:
        html_content = fh.read()

        return convert_html_to_md(html_content)


def DocxConverter(local_path: str):
    """
    Convert a DOCX file to Markdown format.

    Uses mammoth library to first convert DOCX to HTML, then converts
    the HTML to Markdown.

    Args:
        local_path: Path to the DOCX file to convert.

    Returns:
        DocumentConverterResult containing the converted Markdown text.
    """
    with open(local_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        html_content = result.value
    return convert_html_to_md(html_content)


def XlsxConverter(local_path: str):
    """
    Converts Excel files to Markdown using openpyxl.
    Preserves color formatting and other cell styling information.

    Args:
        local_path: Path to the Excel file

    Returns:
        DocumentConverterResult with the Markdown representation of the Excel file
    """
    # Load the workbook
    wb = openpyxl.load_workbook(local_path, data_only=True)
    md_content = ""

    # Helper function to convert RGB color to hex
    def rgb_to_hex(rgb_value):
        if not rgb_value:
            return None

        # Convert RGB value to string for processing
        rgb_string = str(rgb_value)

        # Handle RGB format like 'RGB(255, 255, 255)'
        if isinstance(rgb_value, str) and rgb_string.startswith("RGB"):
            rgb_match = re.match(r"RGB\((\d+), (\d+), (\d+)\)", rgb_string)
            if rgb_match:
                r, g, b = map(int, rgb_match.groups())
                return f"#{r:02x}{g:02x}{b:02x}"

        # Special handling for FFFFFFFF (white) and 00000000 (transparent/none)
        if rgb_string in ["FFFFFFFF", "00000000", "none", "auto"]:
            return None

        # Handle ARGB format (common in openpyxl)
        if len(rgb_string) == 8:  # ARGB format like 'FF5733FF'
            return f"#{rgb_string[2:]}"  # Strip alpha channel

        # Handle direct hex values like 'FF5733'
        if isinstance(rgb_value, str):
            return f"#{rgb_string}" if not rgb_string.startswith("#") else rgb_string

        return None  # Return None for unrecognized formats

    # Helper function to detect and format cell styling
    def get_cell_format_info(cell):
        info = {}

        # Get background color if it exists
        if cell.fill and hasattr(cell.fill, "fgColor") and cell.fill.fgColor:
            # Get the RGB value - in openpyxl this can be stored in different attributes
            rgb_value = None
            if hasattr(cell.fill.fgColor, "rgb") and cell.fill.fgColor.rgb:
                rgb_value = cell.fill.fgColor.rgb
            elif hasattr(cell.fill.fgColor, "value") and cell.fill.fgColor.value:
                rgb_value = cell.fill.fgColor.value

            if rgb_value:
                bg_color = rgb_to_hex(rgb_value)
                if bg_color:  # Skip transparent or white (handled in rgb_to_hex)
                    info["bg_color"] = bg_color

        # Get font color if it exists
        if cell.font and hasattr(cell.font, "color") and cell.font.color:
            # Get the RGB value - in openpyxl this can be stored in different attributes
            rgb_value = None
            if hasattr(cell.font.color, "rgb") and cell.font.color.rgb:
                rgb_value = cell.font.color.rgb
            elif hasattr(cell.font.color, "value") and cell.font.color.value:
                rgb_value = cell.font.color.value

            if rgb_value:
                font_color = rgb_to_hex(rgb_value)
                if font_color:  # Skip transparent (handled in rgb_to_hex)
                    info["font_color"] = font_color

        # Get font weight (bold)
        if cell.font and cell.font.bold:
            info["bold"] = True

        # Get font style (italic)
        if cell.font and cell.font.italic:
            info["italic"] = True

        # Get font underline
        if cell.font and cell.font.underline and cell.font.underline != "none":
            info["underline"] = True

        return info

    # Process each sheet in the workbook
    for sheet_name in wb.sheetnames:
        try:
            sheet = wb[sheet_name]
            md_content += f"## {sheet_name}\n\n"

            # Get the dimensions of the used part of the sheet
            min_row, min_col = 1, 1
            max_row = max(
                (cell.row for cell in sheet._cells.values() if cell.value is not None),
                default=0,
            )
            max_col = max(
                (
                    cell.column
                    for cell in sheet._cells.values()
                    if cell.value is not None
                ),
                default=0,
            )

            if max_row == 0 or max_col == 0:
                md_content += "This sheet is empty.\n\n"
                continue
        except Exception as e:
            error_msg = f"Error processing sheet '{sheet_name}': {str(e)}"
            print(error_msg)
            md_content += (
                f"## {sheet_name}\n\nError processing this sheet: {str(e)}\n\n"
            )
            continue

        try:
            # First, determine column widths
            col_widths = {}
            for col_idx in range(min_col, max_col + 1):
                max_length = 0
                # col_letter = get_column_letter(col_idx)
                _ = get_column_letter(col_idx)
                for row_idx in range(min_row, max_row + 1):
                    try:
                        cell = sheet.cell(row=row_idx, column=col_idx)
                        cell_value = str(cell.value) if cell.value is not None else ""
                        max_length = max(max_length, len(cell_value))
                    except Exception as e:
                        print(
                            f"Warning: Error processing cell at row {row_idx}, column {col_idx}: {str(e)}"
                        )
                        max_length = max(max_length, 10)  # Use reasonable default
                col_widths[col_idx] = max(max_length + 2, 5)  # Min width of 5

            # Start building the table
            # Header row with column separators
            md_content += "|"
            for col_idx in range(min_col, max_col + 1):
                md_content += " " + " " * col_widths[col_idx] + " |"
            md_content += "\n"

            # Separator row
            md_content += "|"
            for col_idx in range(min_col, max_col + 1):
                md_content += ":" + "-" * col_widths[col_idx] + ":|"
            md_content += "\n"

            # Data rows
            for row_idx in range(min_row, max_row + 1):
                md_content += "|"
                for col_idx in range(min_col, max_col + 1):
                    try:
                        cell = sheet.cell(row=row_idx, column=col_idx)
                        cell_value = str(cell.value) if cell.value is not None else ""

                        # Get formatting info
                        try:
                            format_info = get_cell_format_info(cell)
                        except Exception as e:
                            print(
                                f"Warning: Error getting formatting for cell at row {row_idx}, column {col_idx}: {str(e)}"
                            )
                            format_info = {}

                        formatted_value = cell_value

                        # Add HTML-style formatting if needed
                        if format_info:
                            style_parts = []

                            if "bg_color" in format_info:
                                style_parts.append(
                                    f"background-color:{format_info['bg_color']}"
                                )

                            if "font_color" in format_info:
                                style_parts.append(f"color:{format_info['font_color']}")

                            span_attributes = []
                            if style_parts:
                                span_attributes.append(
                                    f'style="{"; ".join(style_parts)}"'
                                )

                            # Format with bold/italic/underline if needed
                            inner_value = cell_value
                            if "bold" in format_info:
                                inner_value = f"<strong>{inner_value}</strong>"
                            if "italic" in format_info:
                                inner_value = f"<em>{inner_value}</em>"
                            if "underline" in format_info:
                                inner_value = f"<u>{inner_value}</u>"

                            # Only add a span if we have style attributes
                            if span_attributes:
                                formatted_value = f"<span {' '.join(span_attributes)}>{inner_value}</span>"
                            else:
                                formatted_value = inner_value

                        # Pad to column width and add to markdown
                        padding = col_widths[col_idx] - len(cell_value)
                        padded_value = " " + formatted_value + " " * (padding + 1)
                        md_content += padded_value + "|"
                    except Exception as e:
                        print(
                            f"Error processing cell at row {row_idx}, column {col_idx}: {str(e)}"
                        )
                        # Add a placeholder for the failed cell
                        padded_value = " [Error] " + " " * (col_widths[col_idx] - 7)
                        md_content += padded_value + " |"

                md_content += "\n"
        except Exception as e:
            error_msg = f"Error generating table for sheet '{sheet_name}': {str(e)}\n{traceback.format_exc()}"
            print(error_msg)
            md_content += f"Error generating table: {str(e)}\n\n"

        # Add formatting legend
        has_formatting = False
        for row_idx in range(min_row, max_row + 1):
            for col_idx in range(min_col, max_col + 1):
                cell = sheet.cell(row=row_idx, column=col_idx)
                if get_cell_format_info(cell):
                    has_formatting = True
                    break
            if has_formatting:
                break

        if has_formatting:
            md_content += "\n### Formatting Information\n"
            md_content += "The table above includes HTML formatting to represent colors and styles from the original Excel file.\n"
            md_content += "This formatting may not display in all Markdown viewers.\n"

        md_content += "\n\n"  # Extra newlines between sheets

    return DocumentConverterResult(
        title=None,
        text_content=md_content.strip(),
    )


def PptxConverter(local_path: str) -> DocumentConverterResult:
    """
    Converts PPTX files to Markdown. Supports headings, tables and images with alt text.

    Args:
        local_path: Path to the PPTX file

    Returns:
        DocumentConverterResult containing the converted Markdown text
    """

    def is_picture(shape):
        """Check if a shape is a picture"""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def is_table(shape):
        """Check if a shape is a table"""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    if not local_path.endswith(".pptx"):
        return DocumentConverterResult(
            title=None,
            text_content=f"Error: Expected .pptx file, got: {local_path}",
        )

    md_content = ""
    presentation = pptx.Presentation(local_path)
    slide_num = 0

    for slide in presentation.slides:
        slide_num += 1
        md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"
        title = slide.shapes.title

        for shape in slide.shapes:
            # Pictures
            if is_picture(shape):
                # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                alt_text = ""
                try:
                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                except Exception:
                    pass
                # A placeholder name
                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                md_content += (
                    "\n!["
                    + (alt_text if alt_text else shape.name)
                    + "]("
                    + filename
                    + ")\n"
                )

            # Tables
            if is_table(shape):
                html_table = "<html><body><table>"
                first_row = True
                for row in shape.table.rows:
                    html_table += "<tr>"
                    for cell in row.cells:
                        if first_row:
                            html_table += "<th>" + html.escape(cell.text) + "</th>"
                        else:
                            html_table += "<td>" + html.escape(cell.text) + "</td>"
                    html_table += "</tr>"
                    first_row = False
                html_table += "</table></body></html>"

                # Note: This would require a separate HTML to Markdown converter function
                # In this version, I'm assuming a convert_html_to_md function exists
                md_content += (
                    "\n" + convert_html_to_md(html_table).text_content.strip() + "\n"
                )

            # Text areas
            elif shape.has_text_frame:
                if shape == title:
                    md_content += "# " + shape.text.lstrip() + "\n"
                else:
                    md_content += shape.text + "\n"

        md_content = md_content.strip()
        if slide.has_notes_slide:
            md_content += "\n\n### Notes:\n"
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                md_content += notes_frame.text
            md_content = md_content.strip()

    return DocumentConverterResult(
        title=None,
        text_content=md_content.strip(),
    )


def ZipConverter(local_path: str, **kwargs):
    """
    Extracts ZIP files to a temporary directory and processes each file according to its extension.
    Returns a combined result of all processed files.
    """
    import zipfile

    temp_dir = tempfile.mkdtemp(prefix="zip_extract_")
    md_content = f"# Extracted from ZIP: {os.path.basename(local_path)}\n\n"

    try:
        with zipfile.ZipFile(local_path, "r") as zip_ref:
            zip_ref.extractall(temp_dir)

        # Get all extracted files
        extracted_files = []
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                rel_path = os.path.relpath(file_path, temp_dir)
                extracted_files.append((file_path, rel_path))

        if not extracted_files:
            md_content += "The ZIP file is empty or contains no files.\n"
        else:
            md_content += f"Total files extracted: {len(extracted_files)}\n\n"

            for file_path, rel_path in extracted_files:
                md_content += f"## File: {rel_path}\n\n"

                # Process each file based on its extension
                file_extension = (
                    file_path.rsplit(".", maxsplit=1)[-1].lower()
                    if "." in file_path
                    else ""
                )
                file_result = None

                try:
                    # Use the same processing logic as process_input
                    if file_extension == "py":
                        with open(file_path, "r", encoding="utf-8") as f:
                            file_result = DocumentConverterResult(
                                title=None, text_content=f.read()
                            )

                    elif file_extension in [
                        "txt",
                        "md",
                        "sh",
                        "yaml",
                        "yml",
                        "toml",
                        "csv",
                    ]:
                        with open(file_path, "r", encoding="utf-8") as f:
                            file_result = DocumentConverterResult(
                                title=None, text_content=f.read()
                            )

                    elif file_extension in ["jsonld", "json"]:
                        with open(file_path, "r", encoding="utf-8") as f:
                            file_result = DocumentConverterResult(
                                title=None,
                                text_content=json.dumps(
                                    json.load(f), ensure_ascii=False, indent=2
                                ),
                            )

                    elif file_extension in ["xlsx", "xls"]:
                        file_result = XlsxConverter(local_path=file_path)

                    elif file_extension == "pdf":
                        file_result = DocumentConverterResult(
                            title=None,
                            text_content=pdfminer.high_level.extract_text(file_path),
                        )

                    elif file_extension in ["docx", "doc"]:
                        file_result = DocxConverter(local_path=file_path)

                    elif file_extension in ["html", "htm"]:
                        file_result = HtmlConverter(local_path=file_path)

                    elif file_extension in ["pptx", "ppt"]:
                        file_result = PptxConverter(local_path=file_path)

                    elif file_extension in IMAGE_EXTENSIONS:
                        # Generate image caption for files in ZIP
                        caption = _generate_image_caption(file_path)
                        md_content += "[Image file]\n\n"
                        md_content += f"> {caption}\n\n"
                        continue

                    elif file_extension in AUDIO_EXTENSIONS:
                        # Generate audio caption for files in ZIP
                        caption = _generate_audio_caption(file_path)
                        md_content += "[Audio file]\n\n"
                        md_content += f"> {caption}\n\n"
                        continue

                    elif file_extension in VIDEO_EXTENSIONS:
                        # Generate video caption for files in ZIP
                        caption = _generate_video_caption(file_path)
                        md_content += "[Video file]\n\n"
                        md_content += f"> {caption}\n\n"
                        continue

                    elif file_extension == "pdb":
                        md_content += "[PDB file - specialized format]\n\n"
                        continue

                    else:
                        # Try MarkItDown as fallback
                        try:
                            md_tool = MarkItDown(enable_plugins=True)
                            file_result = md_tool.convert(file_path)
                        except Exception:
                            md_content += (
                                f"[Unsupported file type: {file_extension}]\n\n"
                            )
                            continue

                    # Add the processed content
                    if file_result and getattr(file_result, "text_content", None):
                        content = file_result.text_content
                        # Limit length for each file
                        max_len = 50_000
                        if len(content) > max_len:
                            content = content[:max_len] + "\n... [Content truncated]"
                        md_content += f"```\n{content}\n```\n\n"

                except Exception as e:
                    md_content += f"[Error processing file: {str(e)}]\n\n"
                    print(f"Warning: Error processing {rel_path} from ZIP: {e}")

    finally:
        # Clean up temporary directory
        try:
            shutil.rmtree(temp_dir)
        except Exception as e:
            print(f"Warning: Could not remove temporary directory {temp_dir}: {e}")

    return DocumentConverterResult(
        title="ZIP Archive Contents", text_content=md_content.strip()
    )
